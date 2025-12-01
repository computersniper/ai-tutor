import os
import json
import re
from dataclasses import dataclass
from typing import List, Dict, Optional, Tuple

from deepseek_client import DeepSeekClient  # Reuse your existing wrapper


# =======================
# 0. Course Materials Knowledge Base (Simple RAG -> Global Knowledge Graph)
# =======================

class CourseKnowledgeBase:
    """
    Lightweight Course Knowledge Base (Revised to "Global Knowledge Graph" mode):
    - Scans course materials (txt/md/pdf/pptx) in a specified directory
    - Still splits into multiple "chunks", but no longer retrieves based on query
    - Organizes all materials by "file dimension" into one large text self.global_context
    - When answering questions, directly provides the entire global_context to the LLM,
      allowing the LLM to decide which sections/materials to use (e.g., chapters 1-10 vs. all)
    """

    def __init__(self, folder: str = "course_materials"):
        self.folder = folder
        # Each chunk: {"path": str, "text": str}
        self.chunks: List[Dict] = []
        # The entire course knowledge graph text
        self.global_context: str = ""
        self._load_folder()

    def _load_folder(self):
        if not os.path.isdir(self.folder):
            print(f"[Knowledge Base] Directory {self.folder} does not exist, creating an empty directory first.")
            os.makedirs(self.folder, exist_ok=True)
            return

        print(f"[Knowledge Base] Loading materials from {self.folder}...")
        for root, _, files in os.walk(self.folder):
            for name in files:
                path = os.path.join(root, name)
                ext = os.path.splitext(name)[1].lower()
                try:
                    if ext in [".txt", ".md"]:
                        text = self._load_text_file(path)
                    elif ext == ".pdf":
                        text = self._load_pdf_file(path)
                    elif ext == ".pptx":
                        text = self._load_pptx_file(path)
                    else:
                        continue

                    text = text.strip()
                    if not text:
                        continue

                    # Split the entire content into multiple chunks (mainly for readability and context continuity)
                    for chunk_text in self._chunk_text(text, chunk_size=700, overlap=150):
                        self.chunks.append(
                            {"path": path, "text": chunk_text}
                        )

                    print(f"[Knowledge Base] Loaded and split: {path} (Current total: {len(self.chunks)} chunks)")
                except Exception as e:
                    print(f"[Knowledge Base] Failed to load {path}: {e}")

        print(f"[Knowledge Base] Generated a total of {len(self.chunks)} text fragments (chunks).")

        # Key: Build the "global course knowledge graph" text here in one go
        self.global_context = self._build_global_context()
        print(f"[Knowledge Base] Global course knowledge graph built, total characters approximately {len(self.global_context)}.")

    def _build_global_context(self) -> str:
        """
        Groups all chunks by "file (lecture slides)" and organizes them into a large knowledge graph text.
        This allows the LLM to see at a glance:
        - What materials are available (roughly corresponding to chapters/topics)
        - The specific content of each material (composed of multiple chunks)
        """
        if not self.chunks:
            return ""

        # First group by path: each file ≈ a chapter/topic
        by_file: Dict[str, List[str]] = {}
        for ch in self.chunks:
            by_file.setdefault(ch["path"], []).append(ch["text"])

        parts: List[str] = []
        parts.append("【Course Knowledge Graph Overview】\n")
        parts.append(
            "Below lists all loaded lecture slides/handouts/exercises etc. for this semester. "
            "Each \"material\" roughly corresponds to a chapter or topic. You can selectively "
            "reference parts of it based on the student's question (e.g., review only chapters 1-10, or review all).\n"
        )

        # List a general table of contents so the LLM knows what "chapters/materials" exist
        for idx, path in enumerate(by_file.keys(), 1):
            fname = os.path.basename(path)
            parts.append(f"{idx}. {fname}")
        parts.append("\n【Detailed Content Grouped by Material】\n")

        # Then write the full text of each material (concatenated from multiple chunks)
        for idx, (path, chunk_texts) in enumerate(by_file.items(), 1):
            fname = os.path.basename(path)
            parts.append(f"\n===== Material {idx}: {fname} =====\n")
            parts.append("\n".join(chunk_texts))

        return "\n".join(parts)

    @staticmethod
    def _load_text_file(path: str) -> str:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    def _load_pdf_file(path: str) -> str:
        try:
            import PyPDF2
        except ImportError:
            print("[Knowledge Base] PyPDF2 required for PDF parsing: pip install PyPDF2")
            return ""
        text = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text() or "")
        return "\n".join(text)

    @staticmethod
    def _load_pptx_file(path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            print("[Knowledge Base] python-pptx required for PPTX parsing: pip install python-pptx")
            return ""
        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_lines.append(shape.text)
            if slide_lines:
                # Add blank lines between slides for easier chunking later
                text_runs.append("\n".join(slide_lines))
        return "\n\n".join(text_runs)

    @staticmethod
    def _chunk_text(text: str, chunk_size: int = 700, overlap: int = 150) -> List[str]:
        """
        Simple sliding window text splitting by character length, supports Chinese and English.
        - chunk_size: approximate characters per segment
        - overlap: overlapping characters between adjacent chunks for information continuity
        """
        chunks = []
        start = 0
        length = len(text)
        while start < length:
            end = min(length, start + chunk_size)
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end == length:
                break
            start = end - overlap  # Step back a bit to create overlap
            if start < 0:
                start = 0
        return chunks

    @staticmethod
    def _tokenize(s: str) -> List[str]:
        """
        (No longer used for retrieval logic, but kept for potential future expansion)
        Splits a string into:
        - English/Numbers: continuous segments like 'quick', 'sort', 'lab', 'cs202'
        - Chinese: character by character, e.g., '快','速','排','序'
        So a phrase like "do the lab for the quick sort section" would be split into:
        ['do', 'the', 'lab', 'for', 'the', 'quick', 'sort', 'section']
        For Chinese "做一下quick sort这一节的lab" would be:
        ['做','一','下','quick','sort','这','一','节','的','lab']
        """
        s = s.lower()
        tokens = re.findall(r"[0-9a-zA-Z]+|[\u4e00-\u9fff]", s)
        return tokens

    def search(self, query: str, top_k: int = 4) -> List[Dict]:
        """
        Old "retrieve chunks by query" interface.
        Not recommended for use now; kept for compatibility. Can be reused if retrieval is needed later.
        """
        if not self.chunks:
            return []

        q_tokens = set(self._tokenize(query))
        scored: List[Tuple[float, Dict]] = []

        for chunk in self.chunks:
            d_tokens = set(self._tokenize(chunk["text"]))
            name_tokens = set(self._tokenize(os.path.basename(chunk["path"])))
            all_tokens = d_tokens | name_tokens

            overlap = len(q_tokens & all_tokens)
            if overlap == 0:
                continue

            base_score = overlap / (len(all_tokens) ** 0.5 + 1e-6)
            filename_boost = len(q_tokens & name_tokens) * 0.5
            score = base_score + filename_boost
            scored.append((score, chunk))

        scored.sort(key=lambda x: x[0], reverse=True)
        return [c for _, c in scored[:top_k]]

    def build_context(self, query: str, top_k: int = 4, max_chars: int = 2400) -> str:
        """
        (Key Modification Point)
        Before: Retrieved based on query, took several high-scoring chunks to form context.
        Now: Directly returns the "entire course knowledge graph" self.global_context,
        completely leaving it to the LLM to decide which parts are relevant to the question.

        For backward compatibility, query/top_k/max_chars parameters are kept but not used here.
        If context length becomes a concern in the future, you can truncate self.global_context here.
        """
        if not self.global_context:
            return ""

        # If you want hard truncation, you can change to:
        # if len(self.global_context) > max_chars:
        #     return self.global_context[:max_chars]
        # As per your request: directly give the entire knowledge graph to the LLM.
        return self.global_context


# =======================
# 1. Router Decision
# =======================

ROUTER_SYSTEM_PROMPT = """
You are the [Question Sorting Teaching Assistant Router] for the "Data Structures and Algorithms" course.

The course scope roughly includes:
Arrays, Linked Lists, Stacks, Queues, Trees & Binary Trees, Heaps, Hash Tables, Graphs, Sorting & Searching,
Recursion, Divide and Conquer, Greedy Algorithms, Dynamic Programming, Complexity Analysis, etc.

You need to read the student's question and classify and route it.

1. Field "type":
   - "concept"      : Concept explanation / Algorithm principle / Proof
   - "code"         : Posted code, hoping for debugging or complexity analysis
   - "assignment"   : Homework / Lab / OJ problem
   - "practice"     : "I want to practice / give me some problems / mock exam" etc. (practice requests)
   - "review"       : "Help me review / summarize / organize knowledge points / key points before exam" etc. (review requests)
   - "logistics"    : Course logistics (exam time, homework deadline, etc.)
   - "out_of_scope" : Beyond the scope of this course

2. Field "difficulty": "easy" | "medium" | "hard"

3. Field "need_human_TA" (Whether a human TA is needed):

   Key distinctions:

   A. Cases that can be handled by AI first (need_human_TA = false):
      - Student wants "explanation, hints, ideas, knowledge organization, review key points", etc.;
      - Homework / lab but the tone is "teach me / explain the idea / give hints", not asking for the full answer;
      - Review summaries (exam review), knowledge maps, etc.

   B. Cases that must be handed over to a real human TA (need_human_TA = true):
      - Clearly involves complete answers to official exam (midterm, final, quiz, test) questions;
      - Student explicitly requests "give me the full answer/full code/do my homework for me";
      - The question is clearly beyond the model's capabilities, and you are not confident in giving a reliable answer.

4. Field "route_to" (Routing target):
   - "ConceptAgent" : Concept explanation / Homework guidance focused on ideas
   - "CodeAgent"    : Code debugging, implementation details related questions
   - "PracticeAgent": Generate practice problems, mock practice
   - "ReviewAgent"  : Review overview, key points before exam, study path suggestions
   - "None"         : No further AI processing needed

   For type = "assignment":
   - If it's just hoping for "explanation/hints/teaching how to do the lab", and not an exam question,
     set need_human_TA = false, and choose "ConceptAgent" or "CodeAgent" based on content.

5. Field "notes_for_TA": Notes for the human TA (brief explanation in Chinese is sufficient).

Requirements:
- Strictly output JSON, no extra text.
Structure as follows:
{
  "type": "...",
  "difficulty": "...",
  "need_human_TA": true or false,
  "route_to": "...",
  "notes_for_TA": "..."
}
"""

CONCEPT_SYSTEM_PROMPT = """
You are the [Concept Explanation Teaching Assistant] for the "Data Structures and Algorithms" course.

Language Rules (Very Important！！！！！！！！！！！！！！！！！！！！！！！！！！！！！！！！！！！):
- Automatically detect the primary language of the student's question;
- If the student asks mainly in **English**, answer completely in English (can occasionally include a little Chinese to explain difficult points);
- If the student asks mainly in **Chinese**, answer completely in Chinese (retain necessary English terms like quick sort, pivot, etc.);
- If the question mixes Chinese and English, answer in the language the student uses more.

Very Important:
- You will receive some [Course Material Fragments] from the teacher's PPT/handouts/practice problems.
- When the expressions, definitions, pseudocode in these materials are not exactly the same as the generic version you remember,
  **Follow the [Course Materials]**, prioritize the terminology, symbols, and pseudocode style from the PPT.

Behavior Guidelines:
1. Answer in Chinese by default; if the student clearly prefers English, you can use English.
2. Suggested answer structure:
   (1) Summarize the core idea in 2~3 sentences;
   (2) Explain step by step in detail;
   (3) Try to combine with examples or lab problems from the [Course Materials];
   (4) Explain time complexity, space complexity (if applicable);
   (5) Remind about common mistakes or misconceptions.
3. For questions that are clearly homework / lab / OJ problems:
   - First, based on the description in the [Course Materials], clearly restate the requirements of the problem/lab;
   - Then provide problem-solving ideas, key steps, connection to in-class examples;
   - Do not provide complete final answers or complete code that can be directly copied.

If the student says "don't give the answer yet / just give problems without solutions / or similar meaning",
- Then only provide the problem list, without reference solutions;
If the student says "give reference answers and explanations after the problems / or similar meaning",
- Then include brief answers and explanations after each problem.

"""

CODE_SYSTEM_PROMPT = """
You are the [Code Debugging Teaching Assistant] for the "Data Structures and Algorithms" course.

Language Rules (Very Important):
- Automatically detect the primary language of the student's question;
- If the student asks mainly in **English**, answer completely in English (can occasionally include a little Chinese to explain difficult points);
- If the student asks mainly in **Chinese**, answer completely in Chinese (retain necessary English terms like quick sort, pivot, etc.);
- If the question mixes Chinese and English, answer in the language the student uses more.

Tasks:
- Analyze the code provided by the student (C/C++/Java/Python, etc.);
- Find logical errors, boundary issues, high complexity, etc.;
- Provide modification suggestions and explanations.

Behavior Guidelines:
1. Suggested answer structure:
   (1) Summarize the code's intent and potential issues;
   (2) Point out suspicious locations or code snippets;
   (3) Explain the cause of the error;
   (4) Provide modification suggestions, can use pseudocode or partial modification examples;
   (5) Analyze time complexity / space complexity, and suggest if there is room for optimization.
2. For homework code, avoid providing the complete final version that can be directly submitted.
3. If there is [Course Materials] context, you can refer to the correct writing style or pseudocode style within.

If the student says "don't give the answer yet / just give problems without solutions / or similar meaning",
- Then only provide the problem list, without reference solutions;
If the student says "give reference answers and explanations after the problems / or similar meaning",
- Then include brief answers and explanations after each problem.

"""

PRACTICE_SYSTEM_PROMPT = """
You are the [Practice Problem Generation Teaching Assistant] for the "Data Structures and Algorithms" course.

Language Rules (Very Important):
- Automatically detect the primary language of the student's question;
- If the student asks mainly in **English**, answer completely in English (can occasionally include a little Chinese to explain difficult points);
- If the student asks mainly in **Chinese**, answer completely in Chinese (retain necessary English terms like quick sort, pivot, etc.);
- If the question mixes Chinese and English, answer in the language the student uses more.

You can generate practice problems based on the course materials to help students master knowledge.
Please adhere to:

1. Design an appropriate number of practice problems according to the student's requirements (chapter, difficulty, number of problems).
2. Question types can include:
   - Conceptual Q&A (short answer);
   - Complexity determination questions;
   - Code reading questions;
   - Algorithm design questions (mainly describe ideas).
3. By default, provide "Problem + Reference Answer/Explanation". If the student explicitly says "don't give answers yet", only provide the problems.
4. Prioritize referencing the [Course Material Fragments] you receive, so the problem style aligns with the real course.
5. If the student doesn't specify difficulty, default to medium.

If the student says "don't give the answer yet / just give problems without solutions / or similar meaning",
- Then only provide the problem list, without reference solutions;
If the student says "give reference answers and explanations after the problems / or similar meaning",
- Then include brief answers and explanations after each problem.
"""

REVIEW_SYSTEM_PROMPT = """
You are the [Review and Summary Teaching Assistant Review TA] for the "Data Structures and Algorithms" course.

Language Rules (Very Important):
- Automatically detect the primary language of the student's question;
- If the student asks mainly in English, answer completely in English;
- If the student asks mainly in Chinese, answer completely in Chinese (retain necessary English terms);
- If it's a mix, answer in the language the student uses more.

Your goal is to help students review this course efficiently, especially:
- Organize knowledge points by chapter or topic (e.g., the chapter on quick sort, all tree-related content);
- Summarize for each knowledge point: core idea, typical operations/algorithms, time & space complexity, common test points;
- Based on course materials, list "typical question types that might appear in homework or exams" (only mention the type, don't give full exam answers);
- Provide simple self-test/practice suggestions (can be a few small questions for students to check themselves).

Behavior Guidelines:
1. Prioritize using the content and terminology from the provided [Course Material Fragments].
2. Suggested answer structure (can be adjusted as needed):
   (1) The overall goal and main line of this part/chapter;
   (2) List of key knowledge points (bullet list);
   (3) For each knowledge point: one-sentence mnemonic + typical operation/algorithm + complexity;
   (4) Common pitfalls/confusion points;
   (5) Suggested review order and practice direction.
3. If the student gives a "review time limit" (e.g., exam in 3 days),
   you can provide a simple study plan and suggested daily tasks.
4. Do not directly leak complete standard answers for exams or quizzes.
"""


# =======================
# Router Data Structure
# =======================

@dataclass
class RouterDecision:
    type: str
    difficulty: str
    need_human_TA: bool
    route_to: str
    notes_for_TA: str


def extract_json_block(text: str) -> str:
    """Extracts the first {...} JSON block from the model output."""
    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1:
        raise ValueError(f"Router output is not valid JSON: {text!r}")
    return text[start: end + 1]


# =======================
# 2. Agent Wrappers
# =======================

class RouterAgent:
    def __init__(self, client: DeepSeekClient):
        self.client = client

    def route(self, question: str) -> RouterDecision:
        resp = self.client.chat(
            system_prompt=ROUTER_SYSTEM_PROMPT,
            user_message=f"The student's question is as follows, please output the decision in the agreed format:\n\n{question}",
            temperature=0.0,
        )
        data = json.loads(extract_json_block(resp))
        return RouterDecision(
            type=data.get("type", "out_of_scope"),
            difficulty=data.get("difficulty", "medium"),
            need_human_TA=bool(data.get("need_human_TA", False)),
            route_to=data.get("route_to", "None"),
            notes_for_TA=data.get("notes_for_TA", ""),
        )


class ConceptAgent:
    def __init__(self, client: DeepSeekClient, kb: CourseKnowledgeBase):
        self.client = client
        self.kb = kb

    def answer(self, question: str, conversation_history: List[Dict] = None) -> str:
        # Now build_context returns the entire course knowledge graph
        context = self.kb.build_context(question, top_k=4, max_chars=2400)

        # Simple check if it's lab / homework explanation
        lower_q = question.lower()
        lab_keywords = ["lab", "exercise", "experiment", "homework", "assignment", "作业", "练习", "实验"]
        is_lab_like = any(k in lower_q for k in lab_keywords)

        prompt_parts = []
        if context:
            prompt_parts.append("Below is the [Complete Knowledge Graph/Course Materials] for this course (grouped by lecture/chapter), please strictly refer to it:\n")
            prompt_parts.append(context)
            prompt_parts.append("\n---\n")
            print("\n[Debug] Current course knowledge graph context passed to LLM (first 500 chars):\n")
            print(context[:500])
            print("\n[Debug] ---- End ----\n")

        if is_lab_like:
            prompt_parts.append(
                "The above likely already contains the description or related examples for this lab/homework.\n"
                "Please first, based on the materials, **restate the lab's task requirements in your own words**, then provide step-by-step problem-solving ideas,"
                "explain the correspondence with class examples/pseudocode, and finally point out common mistakes."
                "Please do not directly provide complete code or final answers that can be copied.\n"
            )

        prompt_parts.append(f"The student's question is:\n{question}\n")

        user_message = "\n".join(prompt_parts)
        
        # If there is history, build the complete conversation context
        if conversation_history:
            messages = conversation_history + [{"role": "user", "content": user_message}]
            return self.client.chat_with_history(
                system_prompt=CONCEPT_SYSTEM_PROMPT,
                messages=messages,
                temperature=0.5,
            )
        else:
            return self.client.chat(
                system_prompt=CONCEPT_SYSTEM_PROMPT,
                user_message=user_message,
                temperature=0.5,
            )


class CodeAgent:
    def __init__(self, client: DeepSeekClient, kb: CourseKnowledgeBase):
        self.client = client
        self.kb = kb

    def answer(self, question: str, conversation_history: List[Dict] = None) -> str:
        # Also pass the entire knowledge graph to help the LLM compare with pseudocode/implementation style in the materials
        context = self.kb.build_context(question, top_k=2, max_chars=1500)
        prompt_parts = []
        if context:
            prompt_parts.append("Below is the [Complete Knowledge Graph/Course Materials] for this course (grouped by lecture/chapter), for reference:\n")
            prompt_parts.append(context)
            prompt_parts.append("\n---\n")
        prompt_parts.append(f"The student's posted code or question is as follows:\n{question}\n\nPlease analyze and answer as a code debugging TA.")

        user_message = "\n".join(prompt_parts)
        
        # If there is history, build the complete conversation context
        if conversation_history:
            messages = conversation_history + [{"role": "user", "content": user_message}]
            return self.client.chat_with_history(
                system_prompt=CODE_SYSTEM_PROMPT,
                messages=messages,
                temperature=0.4,
            )
        else:
            return self.client.chat(
                system_prompt=CODE_SYSTEM_PROMPT,
                user_message=user_message,
                temperature=0.4,
            )


class PracticeAgent:
    def __init__(self, client: DeepSeekClient, kb: CourseKnowledgeBase):
        self.client = client
        self.kb = kb

    def generate(self, request: str, conversation_history: List[Dict] = None) -> str:
        """
        request: The student's practice needs, e.g.:
        - "Want to practice binary tree traversal, 5 medium difficulty problems, with explanations"
        - "Please generate 10 comprehensive multiple-choice questions for the first three chapters of my course, no answers yet"
        """
        # Generate problems directly based on the complete course knowledge graph
        context = self.kb.build_context(request, top_k=5, max_chars=2500)
        prompt_parts = []
        if context:
            prompt_parts.append("Below is the [Complete Knowledge Graph/Course Materials] for this course (grouped by lecture/chapter), please generate problems based on this content:\n")
            prompt_parts.append(context)
            prompt_parts.append("\n---\n")
        prompt_parts.append(f"Student's practice request:\n{request}\n\nPlease generate suitable practice problems (pay attention to follow the rules in the system prompt).")

        user_message = "\n".join(prompt_parts)
        
        # If there is history, build the complete conversation context
        if conversation_history:
            messages = conversation_history + [{"role": "user", "content": user_message}]
            return self.client.chat_with_history(
                system_prompt=PRACTICE_SYSTEM_PROMPT,
                messages=messages,
                temperature=0.6,
            )
        else:
            return self.client.chat(
                system_prompt=PRACTICE_SYSTEM_PROMPT,
                user_message=user_message,
                temperature=0.6,
            )


class ReviewAgent:
    """
    Review / Summary Teaching Assistant:
    - Generates review outlines, key point lists for a certain part based on student request + course materials.
    """

    def __init__(self, client: DeepSeekClient, kb: CourseKnowledgeBase):
        self.client = client
        self.kb = kb

    def review(self, request: str, conversation_history: List[Dict] = None) -> str:
        # No longer "retrieve based on request", but give the entire course materials to the LLM,
        # letting it choose the corresponding parts to summarize based on keywords like "chapters 1-10 / all".
        context = self.kb.build_context(request, top_k=6, max_chars=2600)

        prompt_parts = []
        if context:
            prompt_parts.append("Below is the [Complete Knowledge Graph/Course Materials] for this course (grouped by lecture/chapter), please perform review summary based on this content:\n")
            prompt_parts.append(context)
            prompt_parts.append("\n---\n")

        prompt_parts.append(
            "The student's review needs are as follows (may include chapter names, knowledge points, exam time, etc.):\n"
        )
        prompt_parts.append(request)
        prompt_parts.append(
            "\n\nPlease generate a structured review guide based on the above course materials and needs."
        )

        user_message = "\n".join(prompt_parts)
        
        # If there is history, build the complete conversation context
        if conversation_history:
            messages = conversation_history + [{"role": "user", "content": user_message}]
            return self.client.chat_with_history(
                system_prompt=REVIEW_SYSTEM_PROMPT,
                messages=messages,
                temperature=0.5,
            )
        else:
            return self.client.chat(
                system_prompt=REVIEW_SYSTEM_PROMPT,
                user_message=user_message,
                temperature=0.5,
            )


# =======================
# 3. Main Teaching Assistant Controller (Adds History Functionality)
# =======================

PENDING_FILE = "pending_for_human.jsonl"
HISTORY_FILE = "conversation_history.json"

class TeachingAssistant:
    """
    Unified external "AI Teaching Assistant" entry point.
    - handle_question: Processes a student's question
    - CLI demo: Command line mode simulating student questions
    """

    def __init__(self, materials_folder: str = "course_materials", session_id: str = "default"):
        self.client = DeepSeekClient()
        self.kb = CourseKnowledgeBase(materials_folder)
        self.session_id = session_id
        
        # Add conversation history management
        self.conversation_history: List[Dict] = []
        self.history_file = f"history_{session_id}.json"
        self._load_history()

        self.router = RouterAgent(self.client)
        self.concept_agent = ConceptAgent(self.client, self.kb)
        self.code_agent = CodeAgent(self.client, self.kb)
        self.practice_agent = PracticeAgent(self.client, self.kb)
        self.review_agent = ReviewAgent(self.client, self.kb)

    def _load_history(self):
        """Load conversation history from file"""
        try:
            if os.path.exists(self.history_file):
                with open(self.history_file, "r", encoding="utf-8") as f:
                    self.conversation_history = json.load(f)
                print(f"[History] Loaded {len(self.conversation_history)} history records")
        except Exception as e:
            print(f"[History] Failed to load history: {e}")
            self.conversation_history = []

    def _save_history(self):
        """Save conversation history to file"""
        try:
            with open(self.history_file, "w", encoding="utf-8") as f:
                json.dump(self.conversation_history, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"[History] Failed to save history: {e}")

    def _add_to_history(self, role: str, content: str):
        """Add a message to history"""
        self.conversation_history.append({"role": role, "content": content})
        # Limit history length to avoid too many tokens
        if len(self.conversation_history) > 20:  # Keep last 10 turns of conversation (20 messages)
            self.conversation_history = self.conversation_history[-20:]
        self._save_history()

    def clear_history(self):
        """Clear conversation history"""
        self.conversation_history = []
        self._save_history()
        print("[History] Conversation history cleared")

    def get_history_preview(self) -> str:
        """Get a preview of the history"""
        if not self.conversation_history:
            return "No conversation history yet"
        
        preview = []
        for i, msg in enumerate(self.conversation_history[-5:], 1):  # Show last 5 messages
            role = "Student" if msg["role"] == "user" else "TA"
            content = msg["content"][:50] + "..." if len(msg["content"]) > 50 else msg["content"]
            preview.append(f"{i}. {role}: {content}")
        return "\n".join(preview)

    @staticmethod
    def save_pending(question: str, decision: RouterDecision, ai_answer: Optional[str]):
        record = {
            "question": question,
            "router": decision.__dict__,
            "ai_answer": ai_answer,
        }
        with open(PENDING_FILE, "a", encoding="utf-8") as f:
            f.write(json.dumps(record, ensure_ascii=False) + "\n")

    def handle_question(self, question: str) -> Dict:
        """
        External call:
        - First use Router to determine type, difficulty, need for human;
        - If human needed: record in pending_for_human.jsonl, AI gives no formal answer;
        - Otherwise call the corresponding Agent based on route_to.
        """
        # Add user question to history
        self._add_to_history("user", question)

        decision = self.router.route(question)

        # --- Assignment fallback strategy: Labs/Homework prioritized for AI to explain 'ideas', exams go to human ---
        if decision.type == "assignment":
            lower_q = question.lower()
            exam_keywords = ["exam", "midterm", "final", "quiz", "test", "考试", "期中", "期末", "小测", "考卷", "试卷"]
            is_exam_like = any(k in lower_q for k in exam_keywords)

            if not is_exam_like:
                # Non-exam assignments, default to allowing AI to give idea guidance first
                decision.need_human_TA = False
                # Choose routing: if mentions "code/C++/Python" etc. go to CodeAgent, otherwise ConceptAgent
                code_keywords = ["code", "c++", "java", "python", "implementation", "compile", "error", "代码", "实现", "报错"]
                if any(k in lower_q for k in code_keywords):
                    decision.route_to = "CodeAgent"
                else:
                    decision.route_to = "ConceptAgent"

        result = {
            "router": decision.__dict__,
            "ai_answer": None,
            "message": "",
        }

        # ---- The following logic is similar to before ----
        if decision.need_human_TA:
            self.save_pending(question, decision, ai_answer=None)
            result["message"] = "This question has been marked as requiring a human TA, added to the pending list."
            return result

        # Call the corresponding Agent, pass history (excluding the current question, as it's already in history)
        history_for_agent = self.conversation_history[:-1]  # Excludes current user question
        
        if decision.route_to == "ConceptAgent":
            answer = self.concept_agent.answer(question, history_for_agent)
        elif decision.route_to == "CodeAgent":
            answer = self.code_agent.answer(question, history_for_agent)
        elif decision.route_to == "PracticeAgent":
            answer = self.practice_agent.generate(question, history_for_agent)
        elif decision.route_to == "ReviewAgent":
            answer = self.review_agent.review(question, history_for_agent)
        else:
            answer = (
                "This question is more suitable for direct answering by the course instructor or a human TA, "
                "the AI TA will not provide a specific answer for now."
            )

        # Add AI answer to history
        self._add_to_history("assistant", answer)

        result["ai_answer"] = answer
        return result


# =======================
# 4. Command Line Demo (Adds History Management Commands)
# =======================

def main():
    print("=== Data Structures and Algorithms AI Multi-Agent Teaching Assistant ===")
    print("Hi! I am the AI TA for this course, I can help you with these things:")
    print("1) 📘 Concept explanation")
    print("   - Example: What is stable sort? Why is quick sort O(n log n) on average?")
    print("2) 🧪 Lab / Homework ideas (no full solutions)")
    print("   - Example: Teach me how to do the lab for the quick sort section, give me ideas and steps.")
    print("3) 🧩 Code debugging")
    print("   - Example: This is my binary search code, it always infinite loops, help me find the bug.")
    print("4) 🎯 Practice problem generation")
    print("   - Example: Want to practice AVL tree, 5 medium difficulty problems with explanations.")
    print("5) 📝 Review & summary")
    print("   - Example: Help me review quick sort and merge sort, summarize key points and common test points.")
    print("")
    print("History Features:")
    print("- Type 'history' or '历史' to view recent conversation")
    print("- Type 'clear history' or '清空历史' to clear conversation history")
    print("- Type 'help' or '?' for help")
    print("- Type 'exit' or 'quit' to exit")
    print("")

    # Can set different session_id to distinguish between different students' conversations
    session_id = input("Please enter session ID (press Enter for default session): ").strip()
    if not session_id:
        session_id = "default"
    
    ta = TeachingAssistant(materials_folder="course_materials", session_id=session_id)

    while True:
        try:
            q = input("Student> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nGoodbye.")
            break

        if not q:
            continue
            
        # Handle history-related commands
        if q.lower() in {"history", "历史", "对话历史"}:
            print("\n【Recent Conversation History】")
            print(ta.get_history_preview())
            print("")
            continue
            
        if q.lower() in {"clear history", "清空历史", "clear", "清空"}:
            ta.clear_history()
            print("Conversation history cleared")
            continue
            
        if q.lower() in {"exit", "quit"}:
            print("Goodbye.")
            break
            
        if q.lower() in {"help", "?", "帮助"}:
            print("\nI can help you with these things:")
            print("1) Concept explanation")
            print("2) Lab / assignment guidance (no full solutions)")
            print("3) Code debugging")
            print("4) Practice problem generation")
            print("5) Review & exam prep")
            print("")
            print("History Commands:")
            print("- history / 历史 : View conversation history")
            print("- clear history / 清空历史 : Clear conversation history")
            print("- help / ? : Show this help")
            print("- exit / quit : Exit program\n")
            continue

        try:
            out = ta.handle_question(q)
        except Exception as e:
            print(f"[Error] Processing failed: {e}")
            continue

        print("\n[Router Decision]:")
        print(json.dumps(out["router"], ensure_ascii=False, indent=2))

        if out.get("ai_answer"):
            print("\n[AI TA Answer]:")
            print(out["ai_answer"])
        else:
            print("\n[Note]:", out.get("message", "This question has been marked for human TA processing."))

        print("-" * 60)


if __name__ == "__main__":
    main()